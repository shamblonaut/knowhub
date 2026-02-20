import os
from django.conf import settings

def extract(resource):
    """Returns list of (page_num, text). Empty list for URLs."""
    if resource.resource_type == 'url':
        # Use metadata as fallback text
        text = f"{resource.title} {resource.description} {' '.join(resource.tags or [])}"
        return [(None, text)]

    path = os.path.join(settings.MEDIA_ROOT, resource.file_path or '')
    if not os.path.exists(path):
        return []

    try:
        fmt = resource.file_format
        if fmt == 'pdf':
            import fitz
            doc = fitz.open(path)
            pages = [(i+1, p.get_text()) for i, p in enumerate(doc) if len(p.get_text()) > 30]
            doc.close()
            return pages
        elif fmt == 'ppt':
            from pptx import Presentation
            prs = Presentation(path)
            return [(i+1, '\n'.join(s.text for s in slide.shapes if hasattr(s,'text') and s.text.strip()))
                    for i, slide in enumerate(prs.slides)]
        elif fmt == 'doc':
            from docx import Document
            paras = [p.text for p in Document(path).paragraphs if len(p.text.strip()) > 10]
            return [(None, '\n'.join(paras[i:i+8])) for i in range(0, len(paras), 8)]
        elif fmt == 'image':
            text = f"{resource.title} {resource.description} {' '.join(resource.tags or [])}"
            return [(None, text)]
    except Exception as e:
        print(f'[Extractor] {resource.id}: {e}')
    return []
